#!/usr/bin/env python3
"""docs/live-slides/*.jpg → docs/ring-deck-live.pptx

PowerPoint / Keynote / Google スライドでそのまま開いて全画面で流せる形にする。
各ページは 3200×1800 の画像を全面に敷く。相手の PC に日本語フォントが無くても
崩れない。発表者ノートには、そのページで話す台本をそのまま入れてある
（PowerPoint の発表者ビューに出る）。

使い方: python3 scripts/build-pptx.py
"""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Inches

SLIDES = Path("docs/live-slides")
SCRIPT = Path("docs/live-script.json")
OUT = Path("docs/ring-deck-live.pptx")

W, H = Inches(13.333), Inches(7.5)   # 16:9

blocks = json.load(open(SCRIPT, encoding="utf-8"))
# スライド番号 → その枚で話す台本
notes: dict[str, str] = {}
for b in blocks:
    n = len(b["slide_list"])
    for i, sl in enumerate(b["slide_list"]):
        tc = f'{int(b["start"] // 60)}:{int(b["start"] % 60):02d}'
        part = "" if n == 1 else f"（この{'前半' if i == 0 else '後半'}で送る）"
        notes[sl] = (
            f'BLOCK {b["no"]}　{tc}〜　約{b["sec"]:.0f}秒 / {len(b["text"])}字{part}\n\n'
            f'{b["text"]}\n'
        )

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
blank = prs.slide_layouts[6]          # completely blank

files = sorted(SLIDES.glob("*.jpg"))
assert files, "docs/live-slides/ に画像がありません"
for f in files:
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(str(f), Emu(0), Emu(0), width=W, height=H)
    key = f.stem
    body = notes.get(key)
    if body is None:
        body = f"（{key} 枚目）"
    s.notes_slide.notes_text_frame.text = body

prs.save(OUT)
size = OUT.stat().st_size / 1024 / 1024
print(f"✓ {OUT}（{len(files)} 枚 / {size:.1f}MB / 16:9 13.333×7.5in）")
print(f"  発表者ノート入り: {sum(1 for f in files if f.stem in notes)} 枚")
